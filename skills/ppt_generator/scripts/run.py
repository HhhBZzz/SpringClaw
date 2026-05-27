#!/usr/bin/env python3
from __future__ import annotations

import json
import os
import re
import sys
from pathlib import Path


SKILL_ID = "ppt_generator"


def parse_payload() -> dict:
    raw = sys.argv[1] if len(sys.argv) > 1 else "{}"
    try:
        payload = json.loads(raw)
        return payload if isinstance(payload, dict) else {"text": str(payload)}
    except Exception:
        return {"text": raw}


def is_truthy(value) -> bool:
    if isinstance(value, bool):
        return value
    return str(value or "").strip().lower() in {"1", "true", "yes", "y", "on"}


def workspace_root() -> Path:
    return Path(os.environ.get("SPRINGCLAW_WORKSPACE_ROOT") or os.environ.get("OPENCLAW_WORKSPACE_ROOT") or os.getcwd()).resolve()


def is_inside(path: Path, root: Path) -> bool:
    try:
        path.resolve().relative_to(root.resolve())
        return True
    except ValueError:
        return False


def resolve_workspace_path(raw: str, *, must_exist: bool = False) -> Path:
    root = workspace_root()
    candidate = Path(str(raw).strip()).expanduser()
    if not candidate.is_absolute():
        candidate = root / candidate
    target = candidate.resolve()
    if not is_inside(target, root):
        raise ValueError("path must be inside workspace")
    if must_exist and not target.is_file():
        raise ValueError("input file does not exist")
    return target


def safe_filename(value: str, suffix: str) -> str:
    base = re.sub(r"[^A-Za-z0-9._-]+", "_", str(value or "slides")).strip("._-")
    if not base:
        base = "slides"
    if base.lower().endswith(suffix.lower()):
        base = base[:-len(suffix)]
    base = base.strip("._-") or "slides"
    max_base_len = max(1, 120 - len(suffix))
    return base[:max_base_len] + suffix


def default_output(filename: str) -> Path:
    return workspace_root() / "data" / "skills" / SKILL_ID / safe_filename(filename, ".pptx")


def read_input_text(payload: dict) -> tuple[str, str]:
    if payload.get("text") or payload.get("content"):
        return str(payload.get("text") or payload.get("content")), "inline"
    file_value = payload.get("inputFile") or payload.get("file") or payload.get("path")
    if file_value:
        target = resolve_workspace_path(str(file_value), must_exist=True)
        return target.read_text(encoding="utf-8", errors="replace")[:2_000_000], str(target)
    return str(payload.get("goal") or ""), "goal"


def result(**kwargs) -> None:
    print(json.dumps({"skill": SKILL_ID, **kwargs}, ensure_ascii=False, indent=2))


def normalize_slides(payload: dict, text: str, deck_title: str) -> list[dict]:
    raw_slides = payload.get("slides")
    slides: list[dict] = []
    if isinstance(raw_slides, list):
        for index, item in enumerate(raw_slides, start=1):
            if isinstance(item, dict):
                title = str(item.get("title") or item.get("header") or f"第 {index} 页")
                bullets = item.get("bullets") or item.get("items") or []
                if isinstance(bullets, str):
                    bullets = [bullets]
                slides.append({"title": title, "bullets": [str(b) for b in bullets][:8]})
            else:
                slides.append({"title": f"第 {index} 页", "bullets": [str(item)]})
    if slides:
        return slides[:30]

    cleaned = text.strip()
    if not cleaned:
        return [{"title": deck_title, "bullets": ["暂无内容"]}]
    blocks = [block.strip() for block in re.split(r"\n\s*\n", cleaned) if block.strip()]
    if len(blocks) <= 1:
        lines = [line.strip(" -\t") for line in cleaned.splitlines() if line.strip()]
        if not lines:
            lines = [cleaned]
        if 2 <= len(lines) <= 12:
            for line in lines:
                slides.append({"title": line[:60], "bullets": [line]})
        else:
            chunk_size = 5
            for index in range(0, len(lines), chunk_size):
                chunk = lines[index:index + chunk_size]
                title = chunk[0][:40] if chunk else f"第 {len(slides) + 1} 页"
                slides.append({"title": title, "bullets": chunk[1:] or chunk})
    else:
        for index, block in enumerate(blocks, start=1):
            lines = [line.strip(" -\t") for line in block.splitlines() if line.strip()]
            title = lines[0][:60] if lines else f"第 {index} 页"
            slides.append({"title": title, "bullets": lines[1:] or [block[:160]]})
    return slides[:30]


def generate_pptx(slides: list[dict], title: str, output_path: Path) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        result(
            status="missingDependency",
            missingDependency="python-pptx",
            installHint="python3 -m pip install python-pptx",
            dryRun=False,
        )
        return

    def set_slide_title(slide, text: str) -> None:
        try:
            title_shape = slide.shapes.title
            if title_shape is not None:
                title_shape.text = text
                return
        except Exception:
            pass
        slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.6), Inches(0.7)).text_frame.text = text

    def get_body_frame(slide):
        try:
            if len(slide.placeholders) > 1:
                return slide.placeholders[1].text_frame
        except Exception:
            pass
        return slide.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(8.0), Inches(4.8)).text_frame

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_title(cover, title)
    get_body_frame(cover).text = "Generated by SpringClaw ppt_generator"

    for slide in slides:
        layout = prs.slide_layouts[1]
        page = prs.slides.add_slide(layout)
        set_slide_title(page, str(slide.get("title") or "Untitled")[:120])
        body = get_body_frame(page)
        body.clear()
        bullets = slide.get("bullets") or []
        if not bullets:
            bullets = [" "]
        for index, bullet in enumerate(bullets[:8]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = str(bullet)[:240]
            paragraph.level = 0

    prs.save(str(output_path))
    result(status="success", dryRun=False, outputFile=str(output_path), title=title, slideCount=len(slides) + 1)


def main() -> None:
    payload = parse_payload()
    try:
        text, source = read_input_text(payload)
        title = str(payload.get("title") or "SpringClaw Slides").strip() or "SpringClaw Slides"
        filename = str(payload.get("filename") or title)
        output = resolve_workspace_path(str(payload["outputFile"])) if payload.get("outputFile") else default_output(filename)
        slides = normalize_slides(payload, text, title)
        dry_run = is_truthy(payload.get("dryRun") or payload.get("dry_run"))
        if dry_run:
            result(status="dryRun", dryRun=True, source=source, outputFile=str(output), title=title, slideCount=len(slides) + 1)
            return
        generate_pptx(slides, title, output)
    except Exception as exc:
        result(status="failed", error=str(exc))


if __name__ == "__main__":
    main()
